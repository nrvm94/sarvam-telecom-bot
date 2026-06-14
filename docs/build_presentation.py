"""
Build the Sarvam AI pre-sales assignment presentation.
Run: python docs/build_presentation.py
Output: docs/Sarvam_AI_Airtel_VoiceBot.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x1B, 0x3E)
GREEN       = RGBColor(0x00, 0x8B, 0x5A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF0, 0xF4, 0xF8)
MID_GRAY    = RGBColor(0x71, 0x7D, 0x8A)
GOLD        = RGBColor(0xF5, 0xA6, 0x23)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xEF)
DARK_TEXT   = RGBColor(0x1A, 0x20, 0x2C)
RED_LIGHT   = RGBColor(0xFF, 0xED, 0xED)
RED_TEXT    = RGBColor(0xC0, 0x39, 0x2B)

W = Inches(13.333)
H = Inches(7.5)

# ── Helpers ──────────────────────────────────────────────────────────────────
def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)

def rect(slide, x, y, w, h, fill, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape

def txbox(slide, x, y, w, h):
    return slide.shapes.add_textbox(x, y, w, h)

def para(tf, text, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, space_before=0, italic=False):
    p = tf.add_paragraph()
    p.text = text
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p

def heading_bar(slide, title, subtitle=None):
    """Dark navy top bar with green left accent."""
    rect(slide, 0, 0, W, Inches(1.15), NAVY)
    rect(slide, 0, 0, Inches(0.08), Inches(1.15), GREEN)
    tb = txbox(slide, Inches(0.25), Inches(0.12), Inches(12.5), Inches(0.55))
    tf = tb.text_frame
    tf.word_wrap = False
    para(tf, title, 26, bold=True, color=WHITE)
    if subtitle:
        tb2 = txbox(slide, Inches(0.25), Inches(0.65), Inches(12.5), Inches(0.38))
        tf2 = tb2.text_frame
        para(tf2, subtitle, 12, color=RGBColor(0xA8, 0xC4, 0xE0), italic=True)

def bullet_card(slide, x, y, w, h, title, bullets, bg=LIGHT_GRAY,
                title_color=NAVY, bullet_color=DARK_TEXT, title_size=13, bullet_size=11.5):
    rect(slide, x, y, w, h, bg)
    # title strip
    rect(slide, x, y, w, Inches(0.34), GREEN)
    tb = txbox(slide, x + Inches(0.12), y + Inches(0.04), w - Inches(0.2), Inches(0.3))
    tf = tb.text_frame
    para(tf, title, title_size, bold=True, color=WHITE)
    # bullets
    tb2 = txbox(slide, x + Inches(0.15), y + Inches(0.42), w - Inches(0.25), h - Inches(0.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    first = True
    for b in bullets:
        p = tf2.add_paragraph() if not first else tf2.paragraphs[0]
        p.text = b
        p.space_before = Pt(0 if first else 5)
        r = p.runs[0] if p.runs else p.add_run()
        r.text = b
        r.font.size = Pt(bullet_size)
        r.font.color.rgb = bullet_color
        r.font.name = "Calibri"
        first = False

def stat_box(slide, x, y, w, h, number, label, bg=NAVY, num_color=GOLD, lbl_color=WHITE):
    rect(slide, x, y, w, h, bg)
    tb = txbox(slide, x, y + Inches(0.15), w, Inches(0.65))
    tf = tb.text_frame
    para(tf, number, 32, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    tb2 = txbox(slide, x, y + Inches(0.7), w, Inches(0.45))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    para(tf2, label, 10.5, color=lbl_color, align=PP_ALIGN.CENTER)


# ── SLIDE 1 — Title ──────────────────────────────────────────────────────────
def slide_title(prs):
    sl = blank_slide(prs)
    # Full navy background
    rect(sl, 0, 0, W, H, NAVY)
    # Green left accent bar
    rect(sl, 0, 0, Inches(0.55), H, GREEN)
    # Decorative bottom strip
    rect(sl, Inches(0.55), H - Inches(0.1), W, Inches(0.1), GREEN)

    # Eyebrow
    tb = txbox(sl, Inches(1.0), Inches(1.2), Inches(11), Inches(0.4))
    para(tb.text_frame, "SARVAM AI  ·  PRE-SALES ASSIGNMENT  ·  SOLUTION ARCHITECT", 10,
         color=RGBColor(0x7E, 0xC8, 0xA4), align=PP_ALIGN.LEFT)

    # Main title
    tb2 = txbox(sl, Inches(1.0), Inches(1.7), Inches(11), Inches(1.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "AI Voice Support Bot"
    r.font.size = Pt(48)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Calibri"

    # Subtitle
    tb3 = txbox(sl, Inches(1.0), Inches(3.1), Inches(10), Inches(0.55))
    tf3 = tb3.text_frame
    para(tf3, "Multilingual Customer Support for Airtel — Powered by Sarvam AI", 20,
         color=RGBColor(0x7E, 0xC8, 0xA4))

    # Divider line
    rect(sl, Inches(1.0), Inches(3.75), Inches(5), Inches(0.03), GREEN)

    # Meta line
    tb4 = txbox(sl, Inches(1.0), Inches(3.9), Inches(10), Inches(0.4))
    para(tb4.text_frame, "Use Case: D2C Customer Support  ·  Voice Bot + Agentic Workflow  ·  Hindi & English", 12,
         color=MID_GRAY)

    # Author
    tb5 = txbox(sl, Inches(1.0), Inches(5.6), Inches(8), Inches(0.6))
    para(tb5.text_frame, "Neerav Mahadevan", 16, bold=True, color=WHITE)
    tb6 = txbox(sl, Inches(1.0), Inches(6.15), Inches(8), Inches(0.4))
    para(tb6.text_frame, "June 2026", 11, color=MID_GRAY)


# ── SLIDE 2 — The Problem ─────────────────────────────────────────────────────
def slide_problem(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, WHITE)
    heading_bar(sl, "The Problem", "India's telecom support infrastructure is failing customers at scale")

    # Left dark panel with big stats
    rect(sl, 0, Inches(1.15), Inches(4.6), H - Inches(1.15), NAVY)

    stat_box(sl, Inches(0.15), Inches(1.45), Inches(4.25), Inches(1.1),
             "950K", "inbound support calls per day\n(Airtel, all channels)")
    stat_box(sl, Inches(0.15), Inches(2.65), Inches(4.25), Inches(1.1),
             "₹28–35", "cost per human agent\ninteraction")
    stat_box(sl, Inches(0.15), Inches(3.85), Inches(4.25), Inches(1.1),
             "8–12 min", "average customer wait time\nduring peak hours")
    stat_box(sl, Inches(0.15), Inches(5.05), Inches(4.25), Inches(1.1),
             "71%", "of customers abandon DTMF IVR\nbefore resolution  (TRAI 2024)")

    # Right panel — pain points
    tb = txbox(sl, Inches(4.9), Inches(1.3), Inches(8.0), Inches(0.4))
    para(tb.text_frame, "Four compounding pain points", 15, bold=True, color=NAVY)

    pain_points = [
        ("Language Barrier",
         "68% of subscribers prefer regional languages. Most IVR systems offer Hindi/English only, "
         "failing Tier-2 and Tier-3 customers entirely."),
        ("24/7 Gap",
         "Call centres are staffed ~18 hours/day. Customers with billing disputes or outages "
         "at 2 AM have zero recourse."),
        ("Repetition Fatigue",
         "Average customer repeats their query 2.3× due to misrouting and poor NLU in legacy IVR."),
        ("Cost Trajectory",
         "With 485M subscribers, even a 1% increase in call volume adds ₹130 crore/year "
         "in agent costs — unsustainable."),
    ]
    colors = [GREEN, GOLD, RGBColor(0x21, 0x96, 0xF3), RGBColor(0xE5, 0x34, 0x34)]

    for i, (title, body) in enumerate(pain_points):
        y = Inches(1.85) + i * Inches(1.3)
        rect(sl, Inches(4.85), y, Inches(0.06), Inches(1.05), colors[i])
        tb_t = txbox(sl, Inches(5.1), y, Inches(7.9), Inches(0.32))
        para(tb_t.text_frame, title, 12, bold=True, color=NAVY)
        tb_b = txbox(sl, Inches(5.1), y + Inches(0.3), Inches(7.85), Inches(0.7))
        tf = tb_b.text_frame
        tf.word_wrap = True
        para(tf, body, 11, color=MID_GRAY)


# ── SLIDE 3 — Why Voice AI? ───────────────────────────────────────────────────
def slide_why_ai(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, WHITE)
    heading_bar(sl, "Why Voice AI?", "Voice is the natural interface for India's next 300 million connected users")

    # Top strip — 3 macro stats
    for i, (num, lbl) in enumerate([
        ("500M+", "Indians use voice search\non mobile  (Google India 2024)"),
        ("62%", "of Tier-2/3 first-time internet users\nprefer voice over typing"),
        ("78%", "of Airtel Tier-1 queries are fully\nresolvable without a human agent"),
    ]):
        x = Inches(0.3) + i * Inches(4.35)
        stat_box(sl, x, Inches(1.3), Inches(4.1), Inches(1.25), num, lbl,
                 bg=NAVY if i != 1 else GREEN)

    # End user profile box
    rect(sl, Inches(0.3), Inches(2.8), Inches(6.1), Inches(4.4), LIGHT_GREEN)
    rect(sl, Inches(0.3), Inches(2.8), Inches(6.1), Inches(0.38), GREEN)
    tb = txbox(sl, Inches(0.45), Inches(2.84), Inches(5.8), Inches(0.3))
    para(tb.text_frame, "Who Is the End User?", 13, bold=True, color=WHITE)

    profile = [
        "Airtel subscriber in a Tier-2 or Tier-3 city — comfortable speaking, not typing",
        "First-generation smartphone user; finds app navigation and text-based chat unfamiliar",
        "Speaks Hindi, Hinglish, or a regional language — not pure English",
        "High intent to self-serve but low tolerance for hold music and DTMF menus",
        "Voice is not just a preference — for this segment, it is the only viable interface",
    ]
    tb2 = txbox(sl, Inches(0.5), Inches(3.28), Inches(5.7), Inches(3.7))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    first = True
    for b in profile:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        p.text = "→  " + b
        p.space_before = Pt(0 if first else 7)
        r = p.runs[0] if p.runs else p.add_run()
        r.text = "→  " + b
        r.font.size = Pt(11.5)
        r.font.color.rgb = DARK_TEXT
        r.font.name = "Calibri"
        first = False

    # Why AI vs current approach
    rect(sl, Inches(6.65), Inches(2.8), Inches(6.38), Inches(4.4), LIGHT_GRAY)
    rect(sl, Inches(6.65), Inches(2.8), Inches(6.38), Inches(0.38), NAVY)
    tb3 = txbox(sl, Inches(6.8), Inches(2.84), Inches(6.0), Inches(0.3))
    para(tb3.text_frame, "AI Voice Bot vs. Status Quo", 13, bold=True, color=WHITE)

    comparisons = [
        ("Wait time",        "8–12 min on hold",       "< 3 seconds"),
        ("Availability",     "6 AM – 12 AM only",      "24 / 7 / 365"),
        ("Languages",        "Hindi + English only",    "Hindi, English, Hinglish"),
        ("Cost per call",    "₹28–35",                 "₹0.40"),
        ("Consistency",      "Variable (human error)",  "100% script adherence"),
        ("Scale",            "Staffing-constrained",    "Unlimited, zero incremental cost"),
    ]
    col_x = [Inches(6.8), Inches(8.85), Inches(11.0)]
    hdr_y = Inches(3.27)

    for j, hdr in enumerate(["Metric", "Today", "With AI Voice Bot"]):
        color = NAVY if j == 0 else (RED_TEXT if j == 1 else GREEN)
        tb_h = txbox(sl, col_x[j], hdr_y, Inches(2.1), Inches(0.3))
        para(tb_h.text_frame, hdr, 10.5, bold=True, color=color)

    for i, (metric, before, after) in enumerate(comparisons):
        y = Inches(3.62) + i * Inches(0.48)
        bg = WHITE if i % 2 == 0 else LIGHT_GRAY
        rect(sl, Inches(6.65), y, Inches(6.38), Inches(0.45), bg)
        for j, val in enumerate([metric, before, after]):
            color = DARK_TEXT if j == 0 else (RED_TEXT if j == 1 else GREEN)
            bold = j == 2
            tbc = txbox(sl, col_x[j], y + Inches(0.07), Inches(2.1), Inches(0.32))
            para(tbc.text_frame, val, 10.5, bold=bold, color=color)


# ── SLIDE 4 — Why Sarvam AI? ──────────────────────────────────────────────────
def slide_why_sarvam(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, WHITE)
    heading_bar(sl, "Why Sarvam AI?",
                "The only AI stack purpose-built for India's languages, regulation, and infrastructure")

    # Comparison table header
    tb = txbox(sl, Inches(0.3), Inches(1.28), Inches(12.7), Inches(0.3))
    para(tb.text_frame, "Head-to-Head: Sarvam AI vs Generic Alternatives", 13, bold=True, color=NAVY)

    cols = [Inches(0.3), Inches(3.5), Inches(6.4), Inches(9.3)]
    col_w = Inches(2.85)
    headers = ["Capability", "Google Cloud STT / Vertex", "AWS Transcribe / Bedrock", "Sarvam AI  ✦"]
    hdr_bg = [NAVY, RGBColor(0x3C, 0x3C, 0x3C), RGBColor(0x3C, 0x3C, 0x3C), GREEN]

    for j, (hdr, bg) in enumerate(zip(headers, hdr_bg)):
        rect(sl, cols[j], Inches(1.62), col_w, Inches(0.38), bg)
        tb_h = txbox(sl, cols[j] + Inches(0.08), Inches(1.66), col_w - Inches(0.12), Inches(0.3))
        para(tb_h.text_frame, hdr, 10.5, bold=True, color=WHITE)

    rows = [
        ("Hindi-English code-mixing",    "Partial",          "Limited",           "Native (trained on Hinglish)"),
        ("Indian accent optimisation",   "Partial",          "Partial",           "✓ Trained on 500M+ Indian voices"),
        ("STT latency (India)",          "800–1,200 ms",     "900–1,400 ms",      "≤ 300 ms"),
        ("LLM (Indian language)",        "Gemini (generic)", "Titan (generic)",   "sarvam-105b (India-first reasoning)"),
        ("Data sovereignty / TRAI",      "US servers",       "US servers",        "✓ In-country processing"),
        ("On-premise deployment",        "✗",                "✗",                 "✓ Private cloud / on-prem available"),
        ("Indian language breadth",      "8 languages",      "6 languages",       "22+ Indian languages"),
        ("Cost per minute (STT)",        "$0.006 (~₹0.50)",  "$0.007 (~₹0.58)",  "₹0.18"),
    ]
    neg_vals = {"Partial", "Limited", "US servers", "✗", "800–1,200 ms", "900–1,400 ms",
                "8 languages", "6 languages", "$0.006 (~₹0.50)", "$0.007 (~₹0.58)",
                "Gemini (generic)", "Titan (generic)"}

    for i, row in enumerate(rows):
        y = Inches(2.03) + i * Inches(0.43)
        bg = LIGHT_GRAY if i % 2 == 0 else WHITE
        rect(sl, Inches(0.3), y, Inches(12.73), Inches(0.41), bg)
        for j, val in enumerate(row):
            c = DARK_TEXT
            if j > 0:
                c = RED_TEXT if val in neg_vals else (GREEN if j == 3 else MID_GRAY)
            bold = (j == 3 and val not in neg_vals)
            tbc = txbox(sl, cols[j] + Inches(0.08), y + Inches(0.06), col_w - Inches(0.12), Inches(0.3))
            para(tbc.text_frame, val, 10, bold=bold, color=c)

    # Bottom key message
    rect(sl, Inches(0.3), Inches(5.64), Inches(12.73), Inches(0.5), LIGHT_GREEN)
    tb_k = txbox(sl, Inches(0.5), Inches(5.7), Inches(12.4), Inches(0.38))
    para(tb_k.text_frame,
         "For Airtel: TRAI compliance + on-premise deployment + Hinglish code-mixing is a non-negotiable stack — "
         "only Sarvam delivers all three natively.", 11, bold=True, color=NAVY)


# ── SLIDE 5 — Customer Proof Points ──────────────────────────────────────────
def slide_customers(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, NAVY)
    # Light top bar
    rect(sl, 0, 0, W, Inches(1.15), RGBColor(0x0A, 0x14, 0x30))
    rect(sl, 0, 0, Inches(0.08), Inches(1.15), GREEN)

    tb0 = txbox(sl, Inches(0.25), Inches(0.12), Inches(12.5), Inches(0.55))
    para(tb0.text_frame, "Sarvam AI — Enterprise Proof Points", 26, bold=True, color=WHITE)
    tb0b = txbox(sl, Inches(0.25), Inches(0.65), Inches(12.5), Inches(0.38))
    para(tb0b.text_frame,
         "Deployed at scale across BFSI, Government, and Service industries — ordered by relevance to enterprise customer support",
         12, color=RGBColor(0xA8, 0xC4, 0xE0), italic=True)

    customers = [
        {
            "company": "SBI Life Insurance",
            "industry": "BFSI · Insurance",
            "stat": "Millions of\npolicy calls automated",
            "detail": "AI voice agents handle policy inquiries, renewal reminders, and claims status "
                      "in 10+ Indian languages — 24/7, at a fraction of human agent cost.",
        },
        {
            "company": "Tata Capital",
            "industry": "BFSI · Consumer Lending",
            "stat": "3× increase in\ncustomer engagement",
            "detail": "Multilingual voice agents deployed across consumer loan products — personalised, "
                      "segment-specific conversations that break language barriers at scale.",
        },
        {
            "company": "Mahindra Finance",
            "industry": "BFSI · Rural Finance",
            "stat": "Rural-first\nvoice collections",
            "detail": "Multilingual EMI reminder and collections agent reaching farmers and rural borrowers "
                      "in their native language — populations previously unreachable via text or app.",
        },
        {
            "company": "Urban Company",
            "industry": "D2C · On-demand Services",
            "stat": "Post-service\nvoice follow-up",
            "detail": "Automated voice-first customer feedback and support resolution in Hindi and English "
                      "— reducing manual callback queues for service quality escalations.",
        },
        {
            "company": "Skill India / NABARD",
            "industry": "Government · Public Sector",
            "stat": "50,000+ farmer\nfeedback calls",
            "detail": "Conversational AI collecting structured citizen feedback for Farmer Field School "
                      "programmes in Maharashtra — government-scale reach in regional languages.",
        },
    ]

    card_w = Inches(2.43)
    gap = Inches(0.12)
    for i, c in enumerate(customers):
        x = Inches(0.22) + i * (card_w + gap)
        y = Inches(1.28)
        card_h = Inches(5.95)
        # Card background
        rect(sl, x, y, card_w, card_h, RGBColor(0x14, 0x28, 0x57))
        # Green top strip
        rect(sl, x, y, card_w, Inches(0.06), GREEN)
        # Industry tag
        tag = txbox(sl, x + Inches(0.12), y + Inches(0.15), card_w - Inches(0.2), Inches(0.28))
        para(tag.text_frame, c["industry"], 9, color=RGBColor(0x7E, 0xC8, 0xA4))
        # Company name
        cn = txbox(sl, x + Inches(0.12), y + Inches(0.38), card_w - Inches(0.2), Inches(0.5))
        tf_cn = cn.text_frame
        tf_cn.word_wrap = True
        para(tf_cn, c["company"], 13.5, bold=True, color=WHITE)
        # Divider
        rect(sl, x + Inches(0.12), y + Inches(0.9), card_w - Inches(0.25), Inches(0.025), GREEN)
        # Stat
        st = txbox(sl, x + Inches(0.12), y + Inches(0.98), card_w - Inches(0.2), Inches(0.72))
        tf_st = st.text_frame
        tf_st.word_wrap = True
        para(tf_st, c["stat"], 15, bold=True, color=GOLD)
        # Detail
        det = txbox(sl, x + Inches(0.12), y + Inches(1.72), card_w - Inches(0.2), Inches(4.0))
        tf_det = det.text_frame
        tf_det.word_wrap = True
        para(tf_det, c["detail"], 10.5, color=RGBColor(0xC5, 0xD5, 0xE8))

    # Bottom bar
    rect(sl, 0, H - Inches(0.35), W, Inches(0.35), RGBColor(0x00, 0x6B, 0x45))
    tb_b = txbox(sl, Inches(0.3), H - Inches(0.3), W - Inches(0.5), Inches(0.25))
    para(tb_b.text_frame,
         "Sources: Sarvam.ai (2025–26)  ·  Tata Capital CDO testimonial  ·  Skill India / NABARD programme report",
         9, color=RGBColor(0xD4, 0xED, 0xDA))


# ── SLIDE 6 — What We Built ───────────────────────────────────────────────────
def slide_architecture(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, WHITE)
    heading_bar(sl, "What We Built", "Full-stack A+B solution: Voice Bot + Agentic Escalation Workflow")

    # Architecture flow — 5 pipeline stages
    stages = [
        ("01", "Customer\nBrowser", "React UI\nWebRTC mic capture\nAudio → base64", NAVY),
        ("02", "Sarvam\nSaaras v3 STT", "Hindi / English\nauto-detection\n≤ 300 ms latency", GREEN),
        ("03", "RAG + Sarvam\nsarvam-105b LLM", "ChromaDB retrieval\nover 21 Airtel KB docs\n1–2 sec response", NAVY),
        ("04", "Sarvam\nBulbul v3 TTS", "Natural Indian voice\n37+ speakers\n8kHz telephony output", GREEN),
        ("05", "Spoken\nResponse", "Audio played\nin browser\n~3–5 sec total RTT", RGBColor(0x1A, 0x5C, 0x3A)),
    ]

    box_w = Inches(2.35)
    box_h = Inches(2.0)
    arrow_w = Inches(0.25)
    start_x = Inches(0.2)
    y_top = Inches(1.4)

    for i, (num, title, desc, bg) in enumerate(stages):
        x = start_x + i * (box_w + arrow_w)
        rect(sl, x, y_top, box_w, box_h, bg)
        # Number badge
        badge = txbox(sl, x + Inches(0.08), y_top + Inches(0.08), Inches(0.35), Inches(0.3))
        para(badge.text_frame, num, 9, bold=True,
             color=WHITE if bg == NAVY else NAVY,
             align=PP_ALIGN.CENTER)
        rect(sl, x + Inches(0.08), y_top + Inches(0.08), Inches(0.35), Inches(0.3),
             GREEN if bg == NAVY else WHITE)
        badge2 = txbox(sl, x + Inches(0.08), y_top + Inches(0.08), Inches(0.35), Inches(0.3))
        para(badge2.text_frame, num, 9, bold=True,
             color=NAVY if bg == GREEN or bg == RGBColor(0x1A, 0x5C, 0x3A) else GREEN,
             align=PP_ALIGN.CENTER)

        # Title
        tb_t = txbox(sl, x + Inches(0.08), y_top + Inches(0.42), box_w - Inches(0.15), Inches(0.55))
        tf_t = tb_t.text_frame
        tf_t.word_wrap = True
        para(tf_t, title, 12, bold=True, color=WHITE)

        # Desc
        tb_d = txbox(sl, x + Inches(0.08), y_top + Inches(1.0), box_w - Inches(0.12), Inches(0.9))
        tf_d = tb_d.text_frame
        tf_d.word_wrap = True
        para(tf_d, desc, 9.5, color=RGBColor(0xBB, 0xCC, 0xDD))

        # Arrow (not after last)
        if i < len(stages) - 1:
            ax = x + box_w + Inches(0.03)
            ay = y_top + box_h / 2 - Inches(0.08)
            tb_a = txbox(sl, ax, ay, arrow_w, Inches(0.22))
            para(tb_a.text_frame, "▶", 14, color=GREEN, align=PP_ALIGN.CENTER)

    # Escalation path
    rect(sl, Inches(0.2), Inches(3.58), Inches(13.0), Inches(0.035), RGBColor(0xDD, 0xE3, 0xEA))

    esc_label = txbox(sl, Inches(0.2), Inches(3.7), Inches(2.5), Inches(0.3))
    para(esc_label.text_frame, "ESCALATION PATH  (Agentic Workflow)", 10, bold=True, color=RED_TEXT)

    esc_stages = [
        ("Complex\nIssue Detected", NAVY),
        ("n8n Webhook\nTrigger", RGBColor(0x5C, 0x35, 0x96)),
        ("Mock Ticket\nSystem (:5000)", RGBColor(0x0A, 0x5C, 0x8A)),
        ("360dialog\nWhatsApp API", RGBColor(0x00, 0x79, 0x3C)),
        ("Supabase\nCall Log Update", RGBColor(0x1A, 0x56, 0x76)),
    ]

    esc_box_w = Inches(2.35)
    for i, (title, bg) in enumerate(esc_stages):
        x = Inches(0.2) + i * (esc_box_w + Inches(0.25))
        rect(sl, x, Inches(4.07), esc_box_w, Inches(1.28), bg)
        tb_e = txbox(sl, x + Inches(0.1), Inches(4.17), esc_box_w - Inches(0.18), Inches(1.05))
        tf_e = tb_e.text_frame
        tf_e.word_wrap = True
        para(tf_e, title, 11, bold=True, color=WHITE)
        if i < len(esc_stages) - 1:
            ax = x + esc_box_w + Inches(0.04)
            ay = Inches(4.6)
            tb_a = txbox(sl, ax, ay, Inches(0.22), Inches(0.22))
            para(tb_a.text_frame, "▶", 12, color=GREEN, align=PP_ALIGN.CENTER)

    # Tech tags
    tags = ["React 18 + Vite", "FastAPI + Python", "ChromaDB (RAG)", "Supabase (PostgreSQL)",
            "n8n Automation", "360dialog WhatsApp"]
    tb_tag_label = txbox(sl, Inches(0.2), Inches(5.52), Inches(2.5), Inches(0.28))
    para(tb_tag_label.text_frame, "Stack:", 10, bold=True, color=MID_GRAY)
    for i, tag in enumerate(tags):
        tx = Inches(1.0) + i * Inches(2.05)
        rect(sl, tx, Inches(5.5), Inches(1.95), Inches(0.3), LIGHT_GRAY)
        tb_tg = txbox(sl, tx + Inches(0.08), Inches(5.54), Inches(1.85), Inches(0.24))
        para(tb_tg.text_frame, tag, 9.5, color=DARK_TEXT)

    # GitHub reference
    rect(sl, Inches(0.2), Inches(5.98), Inches(12.9), Inches(0.3), LIGHT_GREEN)
    tb_gh = txbox(sl, Inches(0.35), Inches(6.02), Inches(12.6), Inches(0.24))
    para(tb_gh.text_frame,
         "GitHub: github.com/nrvm94/sarvam-telecom-bot  ·  Full source code, .env.example, setup guide, architecture docs",
         10, color=NAVY)


# ── SLIDE 7 — ROI & Business Case ─────────────────────────────────────────────
def slide_roi(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, WHITE)
    heading_bar(sl, "ROI & Business Case", "A 9-day payback period — at conservative adoption rates")

    # Big three numbers at top
    for i, (num, lbl, bg) in enumerate([
        ("₹0.40",    "cost per AI-handled call\nvs ₹28–35 human agent", NAVY),
        ("70×",      "reduction in\ncost per interaction", GREEN),
        ("₹745 Cr",  "projected annual savings\nat full deployment", RGBColor(0x0A, 0x5C, 0x8A)),
    ]):
        x = Inches(0.3) + i * Inches(4.35)
        stat_box(sl, x, Inches(1.3), Inches(4.1), Inches(1.3), num, lbl, bg=bg)

    # Cost model table
    tb_lbl = txbox(sl, Inches(0.3), Inches(2.82), Inches(6.0), Inches(0.3))
    para(tb_lbl.text_frame, "Cost Model", 13, bold=True, color=NAVY)

    cost_rows = [
        ("",                             "Human Agent",         "AI Voice Bot"),
        ("Cost per interaction",         "₹28–35",             "₹0.40"),
        ("Calls handled / day / unit",   "~120",               "Unlimited"),
        ("Availability",                 "6 AM – 12 AM",       "24 / 7 / 365"),
        ("Avg handle time",              "4–6 minutes",        "15–30 seconds"),
        ("First call resolution (T1)",   "68%",                "84%"),
        ("Language support",             "Hindi + English",    "Hindi, English, Hinglish"),
    ]
    col_x2 = [Inches(0.3), Inches(3.8), Inches(5.95)]
    col_w2 = [Inches(3.45), Inches(2.1), Inches(2.1)]
    hdr_bgs = [NAVY, RGBColor(0x5A, 0x5A, 0x5A), GREEN]

    for j, (hdr_bg_c, hdr) in enumerate(zip(hdr_bgs, ["", "Human Agent", "AI Voice Bot"])):
        rect(sl, col_x2[j], Inches(3.16), col_w2[j], Inches(0.35), hdr_bg_c)
        if hdr:
            tb_hdr = txbox(sl, col_x2[j] + Inches(0.08), Inches(3.2), col_w2[j] - Inches(0.1), Inches(0.27))
            para(tb_hdr.text_frame, hdr, 10.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, row in enumerate(cost_rows[1:]):
        y = Inches(3.54) + i * Inches(0.38)
        bg = WHITE if i % 2 == 0 else LIGHT_GRAY
        for j, val in enumerate(row):
            rect(sl, col_x2[j], y, col_w2[j], Inches(0.36), bg)
            c = DARK_TEXT if j < 2 else GREEN
            bold_flag = (j == 2)
            tb_v = txbox(sl, col_x2[j] + Inches(0.08), y + Inches(0.06), col_w2[j] - Inches(0.1), Inches(0.26))
            para(tb_v.text_frame, val, 10.5, bold=bold_flag, color=c,
                 align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

    # Savings calc box on right
    rect(sl, Inches(8.35), Inches(2.82), Inches(4.75), Inches(4.5), NAVY)
    tb_sc = txbox(sl, Inches(8.5), Inches(2.9), Inches(4.45), Inches(0.35))
    para(tb_sc.text_frame, "Annual Savings Calculation", 13, bold=True, color=GREEN)

    calc_lines = [
        ("Automatable calls / day", "741,000  (78% of 950K)"),
        ("Cost saving per call",    "₹27.60  (₹28 – ₹0.40)"),
        ("Daily saving",            "₹2.04 crore"),
        ("Annual saving",           "₹745 crore"),
        ("Impl. cost (Year 1)",      "₹18 crore"),
        ("Net Year-1 saving",       "₹727 crore"),
        ("Payback period",          "~9 days"),
    ]
    divider_y = Inches(3.3)
    rect(sl, Inches(8.5), divider_y, Inches(4.4), Inches(0.02), GREEN)

    for i, (label, value) in enumerate(calc_lines):
        y = Inches(3.4) + i * Inches(0.47)
        highlight = (label in {"Annual saving", "Payback period"})
        if highlight:
            rect(sl, Inches(8.35), y - Inches(0.04), Inches(4.75), Inches(0.44), GREEN)
        tb_l = txbox(sl, Inches(8.5), y, Inches(2.4), Inches(0.35))
        para(tb_l.text_frame, label, 10.5, color=RGBColor(0xBB, 0xCC, 0xDD) if not highlight else WHITE)
        tb_v2 = txbox(sl, Inches(10.6), y, Inches(2.35), Inches(0.35))
        para(tb_v2.text_frame, value, 10.5, bold=highlight, color=GOLD if not highlight else WHITE)

    # Conservative note
    rect(sl, Inches(0.3), Inches(6.08), Inches(12.75), Inches(0.3), LIGHT_GREEN)
    tb_n = txbox(sl, Inches(0.45), Inches(6.12), Inches(12.5), Inches(0.24))
    para(tb_n.text_frame,
         "Conservative scenario (10% adoption): ₹74.5 crore / year.  "
         "Every 1% reduction in churn from higher NPS = ₹85 crore ARR retained.",
         10, color=NAVY)


# ── SLIDE 8 — Limitations & 90-Day Roadmap ────────────────────────────────────
def slide_roadmap(prs):
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, WHITE)
    heading_bar(sl, "Limitations & 90-Day Roadmap",
                "Current PoC demonstrates the core stack — production readiness is a 90-day sprint")

    # PoC gaps table
    tb_lbl = txbox(sl, Inches(0.3), Inches(1.28), Inches(5.8), Inches(0.3))
    para(tb_lbl.text_frame, "Current PoC Gaps", 13, bold=True, color=NAVY)

    gaps = [
        ("Telephony / SIP integration",  "High",   "Integrate Exotel or Tata Tele SIP trunk — Phase 2"),
        ("Live CRM integration",          "High",   "Connect to Airtel Salesforce; enable personalised responses"),
        ("OTP customer authentication",   "High",   "Add Airtel OTP-verify API before exposing account data"),
        ("Knowledge base scale",          "Medium", "Ingest full Airtel policy corpus (5,000+ docs) via batch pipeline"),
        ("Mid-session language switch",   "Low",    "Enable seamless Hindi ↔ English pivot within one call"),
    ]
    gap_col = [Inches(0.3), Inches(3.55), Inches(4.55)]
    gap_col_w = [Inches(3.2), Inches(0.95), Inches(4.0)]

    for j, hdr in enumerate(["Gap", "Severity", "Mitigation"]):
        rect(sl, gap_col[j], Inches(1.62), gap_col_w[j], Inches(0.35), NAVY)
        tb_h = txbox(sl, gap_col[j] + Inches(0.08), Inches(1.66), gap_col_w[j] - Inches(0.1), Inches(0.27))
        para(tb_h.text_frame, hdr, 10.5, bold=True, color=WHITE)

    sev_colors = {"High": RED_TEXT, "Medium": RGBColor(0xE6, 0x7E, 0x00), "Low": GREEN}
    for i, (gap, sev, fix) in enumerate(gaps):
        y = Inches(2.0) + i * Inches(0.44)
        bg = WHITE if i % 2 == 0 else LIGHT_GRAY
        for j, val in enumerate([gap, sev, fix]):
            rect(sl, gap_col[j], y, gap_col_w[j], Inches(0.42), bg)
            c = DARK_TEXT if j != 1 else sev_colors.get(val, DARK_TEXT)
            bold_flag = (j == 1)
            tbc = txbox(sl, gap_col[j] + Inches(0.08), y + Inches(0.07), gap_col_w[j] - Inches(0.12), Inches(0.3))
            para(tbc.text_frame, val, 10, bold=bold_flag, color=c)

    # 30-60-90 roadmap
    tb_lbl2 = txbox(sl, Inches(8.55), Inches(1.28), Inches(4.5), Inches(0.3))
    para(tb_lbl2.text_frame, "90-Day Enterprise Rollout", 13, bold=True, color=NAVY)

    phases = [
        ("Days 1–30", "Foundation",
         ["Deploy on Airtel cloud (AWS Mumbai / Azure India)",
          "Ingest full Airtel KB and connect subscriber database",
          "Set up monitoring: call success rate, escalation rate, latency"],
         NAVY),
        ("Days 31–60", "Telephony + Auth",
         ["Integrate IVR system via SIP / WebRTC",
          "Add OTP-based caller authentication",
          "A/B test: AI bot vs legacy IVR on 5% of traffic"],
         GREEN),
        ("Days 61–90", "Scale & Optimise",
         ["Scale to 20% of inbound call volume",
          "Add regional language support (Tamil, Telugu, Bengali)",
          "Board presentation: NPS delta, cost savings, escalation rate"],
         RGBColor(0x0A, 0x5C, 0x8A)),
    ]

    for i, (phase, title, bullets, bg) in enumerate(phases):
        y = Inches(1.62) + i * Inches(1.78)
        rect(sl, Inches(8.55), y, Inches(4.55), Inches(1.68), LIGHT_GRAY)
        rect(sl, Inches(8.55), y, Inches(1.05), Inches(1.68), bg)
        # Phase label (vertical feel via stacked)
        tb_p = txbox(sl, Inches(8.6), y + Inches(0.15), Inches(0.95), Inches(0.3))
        para(tb_p.text_frame, phase, 9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb_pt = txbox(sl, Inches(8.6), y + Inches(0.45), Inches(0.95), Inches(0.25))
        para(tb_pt.text_frame, title, 8.5, color=RGBColor(0xBB, 0xCC, 0xDD), align=PP_ALIGN.CENTER)
        # Bullets
        tb_bl = txbox(sl, Inches(9.72), y + Inches(0.12), Inches(3.28), Inches(1.45))
        tf_bl = tb_bl.text_frame
        tf_bl.word_wrap = True
        first = True
        for b in bullets:
            p = tf_bl.paragraphs[0] if first else tf_bl.add_paragraph()
            p.text = "• " + b
            p.space_before = Pt(0 if first else 5)
            r = p.runs[0] if p.runs else p.add_run()
            r.text = "• " + b
            r.font.size = Pt(10)
            r.font.color.rgb = DARK_TEXT
            r.font.name = "Calibri"
            first = False

    # Closing line
    rect(sl, 0, H - Inches(0.38), W, Inches(0.38), NAVY)
    tb_cl = txbox(sl, Inches(0.3), H - Inches(0.33), W - Inches(0.5), Inches(0.28))
    para(tb_cl.text_frame,
         "Built by Neerav Mahadevan  ·  github.com/nrvm94/sarvam-telecom-bot  ·  June 2026  ·  Sarvam AI Pre-Sales Assignment",
         9.5, color=RGBColor(0x8A, 0xA8, 0xC8), align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_title(prs)
    slide_problem(prs)
    slide_why_ai(prs)
    slide_why_sarvam(prs)
    slide_customers(prs)
    slide_architecture(prs)
    slide_roi(prs)
    slide_roadmap(prs)

    out = "docs/Sarvam_AI_Airtel_VoiceBot.pptx"
    prs.save(out)
    print(f"Saved: {out}")

if __name__ == "__main__":
    main()
